# app.py

import os
from flask import Flask, render_template, request, send_file, redirect, url_for
from pytube import YouTube
from moviepy.editor import AudioFileClip
import speech_recognition as sr
from fpdf import FPDF
from pptx import Presentation

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static'

# Function to download audio and convert to text
def youtube_to_text(youtube_url):
    # Download audio
    yt = YouTube(youtube_url)
    audio_stream = yt.streams.filter(only_audio=True).first()
    audio_file_path = audio_stream.download(filename='audio.mp4')
    
    # Convert audio to text
    recognizer = sr.Recognizer()
    with AudioFileClip(audio_file_path) as audio_clip:
        audio_clip.write_audiofile("audio.wav")
    
    with sr.AudioFile("audio.wav") as source:
        audio_data = recognizer.record(source)
        text = recognizer.recognize_google(audio_data)
    
    # Cleanup audio files
    os.remove("audio.mp4")
    os.remove("audio.wav")
    return text

# Function to create a PDF file from text
def create_pdf(text, filename):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    
    for line in text.splitlines():
        pdf.cell(200, 10, txt=line, ln=True)
    
    pdf.output(filename)

# Function to create a PPT file from text
def create_ppt(text, filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    
    for line in text.splitlines():
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = line[:50]  # Set title text to first 50 characters
        content = slide.placeholders[1]
        content.text = line  # Full text in content box
    
    prs.save(filename)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/process', methods=['POST'])
def process():
    youtube_url = request.form['youtubeUrl']
    text = youtube_to_text(youtube_url)

    # Generate PDF and PPT files
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], 'sample.pdf')
    ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], 'sample.pptx')
    create_pdf(text, pdf_path)
    create_ppt(text, ppt_path)

    return render_template(
        'result.html',
        text=text,
        pdf_path=pdf_path,
        ppt_path=ppt_path
    )

@app.route('/download/<path:filename>')
def download(filename):
    return send_file(filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
